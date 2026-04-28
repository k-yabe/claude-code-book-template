#!/usr/bin/env python3
"""AKKODiS ブランド準拠 PPTX を生成する。

設計方針:
- テンプレ (.pptx) には複数の slide_master が含まれており、合計 30+ の
  layout がスライド種別ごと（表紙・アジェンダ・本文・KPI・クロージング）に
  専用設計されている。
- スライド種別ごとに **全 slide_master を横断して最適な layout を自動選定**:
  - 表紙           → 'Wave Landscape' / 'Mesh Yellow & Blue'
  - アジェンダ     → 'Agenda - Computer' (6 slot) / 'Agenda - Mesh Y&B' (8 slot)
  - 本文           → 'Simple Title, Subtitle & Content' / 'Title, Subtitle and one Paragraph'
  - KPI            → 'Six Text Boxes' (6 slot) / 'Centered Text Box 2' (1 slot)
  - クロージング   → 'Thank you' / 'End with Voice'
- 各 layout の placeholder（Res_Title_01 / Res_SubTitle / Res_Paraph_01 /
  Res_Chapter_0N / Res_Num_0N など）にテキストを流し込むだけで、
  デザイナー設計のグラフィックがそのまま活きる。
- IMG_Back の strip や半透明オーバーレイは一切使わない。
- 全テキストに表記補正 (notation.correct) を適用。

入力 (Markdown ライク):
    # タイトル: ...
    # サブタイトル: ...
    # 用途: external | internal      （省略時 external）
    # トーン: dark | white            （省略時 white）
    # 作成: ...

    ## アジェンダ                                ← 明示で自動アジェンダ抑制
    ## セクション名                              ← 通常スライド
    - 箇条書き
    > 発表者ノート

    ## KPI: 指標A=120% | 指標B=95% | 指標C=80%   ← KPI スライド
    ## まとめ                                    ← クロージング

依存: python-pptx
"""

from __future__ import annotations

import argparse
import re
import sys
from dataclasses import dataclass, field
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.util import Pt

sys.path.insert(0, str(Path(__file__).resolve().parent))
from notation import correct  # noqa: E402

NAVY = RGBColor(0x00, 0x1F, 0x33)
GOLD = RGBColor(0xFF, 0xB8, 0x1C)


@dataclass
class Section:
    title: str
    bullets: list[str] = field(default_factory=list)
    notes: list[str] = field(default_factory=list)
    kind: str = "content"  # content | kpi | closing | agenda
    kpis: list[tuple[str, str]] = field(default_factory=list)


@dataclass
class Deck:
    title: str = "Untitled"
    subtitle: str = ""
    audience: str = "external"
    tone: str = "white"
    author: str = ""
    sections: list[Section] = field(default_factory=list)


# ── パース ──────────────────────────────────────────────────


def parse_input(text: str) -> Deck:
    deck = Deck()
    current: Section | None = None
    meta_re = {
        r"^#\s*(?:タイトル|title)\s*[:：]\s*(.+)$": "title",
        r"^#\s*(?:サブタイトル|subtitle)\s*[:：]\s*(.+)$": "subtitle",
        r"^#\s*(?:用途|audience)\s*[:：]\s*(\S+)": "audience",
        r"^#\s*(?:トーン|tone)\s*[:：]\s*(\S+)": "tone",
        r"^#\s*(?:作成|author)\s*[:：]\s*(.+)$": "author",
    }

    def parse_kpi_header(line: str) -> list[tuple[str, str]] | None:
        m = re.match(r"^##\s*KPI\s*[:：]\s*(.+)$", line, re.IGNORECASE)
        if not m:
            return None
        items = []
        for chunk in m.group(1).split("|"):
            if "=" in chunk:
                k, v = chunk.split("=", 1)
                items.append((k.strip(), v.strip()))
        return items

    for raw in text.splitlines():
        line = raw.rstrip()
        if not line.strip() and current is None:
            continue
        if current is None:
            matched = False
            for pattern, attr in meta_re.items():
                m = re.match(pattern, line, re.IGNORECASE)
                if m:
                    val = m.group(1).strip()
                    setattr(deck, attr, val.lower() if attr in ("audience", "tone") else val)
                    matched = True
                    break
            if matched:
                continue
        kpis = parse_kpi_header(line)
        if kpis is not None:
            current = Section(title="KPI ハイライト", kind="kpi", kpis=kpis)
            deck.sections.append(current)
            continue
        m = re.match(r"^##\s+(.+)$", line)
        if m:
            title = m.group(1).strip()
            if title in ("アジェンダ", "Agenda", "目次"):
                kind = "agenda"
            elif title in ("まとめ", "結論", "Conclusion", "Closing", "Thank you"):
                kind = "closing"
            else:
                kind = "content"
            current = Section(title=title, kind=kind)
            deck.sections.append(current)
            continue
        m = re.match(r"^\s*>\s+(.+)$", line)
        if m and current is not None:
            current.notes.append(m.group(1).strip())
            continue
        m = re.match(r"^\s*[-*0-9]+\.?\s+(.+)$", line)
        if m and current is not None:
            current.bullets.append(m.group(1).strip())
            continue
        if current is not None and line.strip():
            current.bullets.append(line.strip())
    return deck


# ── 全 master を横断して layout を探す ─────────────────────────


def _all_layouts(prs: Presentation):
    """全 slide_master の layout を順に yield する。"""
    for master in prs.slide_masters:
        for L in master.slide_layouts:
            yield L


def _find_layout(prs: Presentation, *keywords: str):
    """名前のキーワード一致で layout を探す（最初の一致を返す）。"""
    for L in _all_layouts(prs):
        name = L.name.lower()
        for kw in keywords:
            if kw.lower() in name:
                return L
    return None


def _select_layouts(prs: Presentation, *, n_kpi: int = 0) -> dict:
    """スライド種別ごとに最適な layout を選定。"""
    return {
        "title": _find_layout(
            prs,
            "Wave Landscape",
            "Mesh Yellow & Blue",
            "Mesh & Blue",
        ) or prs.slide_layouts[0],
        "agenda": _find_layout(
            prs,
            "Agenda - Mesh Y&B",
            "Agenda - Computer",
            "Agenda - With Picture",
        ) or _find_layout(prs, "Title, Subtitle and one Paragrah", "Simple Title, Subtitle & Content"),
        # 本文用は内容に応じて切り替えるため複数候補を保持
        "content_paragraph": _find_layout(
            prs,
            "Simple Title, Subtitle & Content",
            "Title, Subtitle and one Paragrah",
        ),
        "content_six": _find_layout(prs, "Six Text Boxes"),
        "content_two_paragraph": _find_layout(prs, "Two Paragraphs with Blue Line"),
        "content_centered": _find_layout(prs, "Centered Text Box 2", "Simple texte centered"),
        "kpi": (
            _find_layout(prs, "Six Text Boxes")
            if n_kpi >= 4
            else _find_layout(prs, "Centered Text Box 2", "Simple texte centered")
        ) or _find_layout(prs, "Simple Title, Subtitle & Content"),
        "closing": _find_layout(prs, "Thank you", "End with Voice", "Eye Akkodis"),
    }


def _pick_content_layout(layouts: dict, section: Section):
    """セクション内容に応じて最適な本文 layout を選ぶ。

    - 4-6 bullets, 各 80 文字以内 → Six Text Boxes (グリッド)
    - 2 bullets, 各 100 文字以内 → Two Paragraphs with Blue Line (2列対比)
    - 1 bullet（長い）または 7+ → Title + Paragraph（番号付き強化）
    """
    bullets = section.bullets or []
    n = len(bullets)
    if n == 0:
        return layouts["content_paragraph"], "paragraph"
    avg_len = sum(len(b) for b in bullets) / n
    max_len = max(len(b) for b in bullets) if bullets else 0

    if 3 <= n <= 6 and max_len <= 80 and layouts.get("content_six") is not None:
        return layouts["content_six"], "six"
    if n == 2 and max_len <= 120 and layouts.get("content_two_paragraph") is not None:
        return layouts["content_two_paragraph"], "two"
    return layouts["content_paragraph"], "paragraph"


# ── placeholder 操作ヘルパ ────────────────────────────────


def _hide_shape(shape) -> None:
    sp = shape._element
    sp.getparent().remove(sp)


def _layout_idx_by_name(layout, *name_parts: str) -> int | None:
    """layout の placeholder で名前部分一致するものの idx を返す。

    keyword は順位通りに優先評価する（先頭ほど優先）。
    """
    # 各 keyword を順に走査し、最初にマッチした placeholder を採用
    for kw in name_parts:
        kw_l = kw.lower()
        for ph in layout.placeholders:
            ph_name = (ph.name or "").lower()
            if kw_l in ph_name:
                return ph.placeholder_format.idx
    return None


def _layout_idxs_by_name(layout, name_part: str) -> list[int]:
    """layout の placeholder で名前部分一致する全 idx を昇順で返す。"""
    result = []
    for ph in layout.placeholders:
        if name_part.lower() in (ph.name or "").lower():
            result.append(ph.placeholder_format.idx)
    return sorted(result)


def _slide_ph_by_idx(slide, idx: int):
    """slide の placeholder を idx で取得。"""
    if idx is None:
        return None
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    return None


def _ph_by_name(slide, *name_parts: str):
    """slide の placeholder を「layout 上の名前」部分一致で探す。

    slide 作成時に placeholder の name は generic 化される（"Text Placeholder N"）が、
    idx は layout から継承されるので、layout の name で逆引きする。
    """
    layout = slide.slide_layout
    idx = _layout_idx_by_name(layout, *name_parts)
    return _slide_ph_by_idx(slide, idx)


def _ph_all_by_name(slide, name_part: str) -> list:
    """layout の name で複数 placeholder を idx 順に取得。"""
    layout = slide.slide_layout
    idxs = _layout_idxs_by_name(layout, name_part)
    return [ph for idx in idxs for ph in [_slide_ph_by_idx(slide, idx)] if ph is not None]


def _set_ph_text(ph, text: str) -> None:
    """placeholder のテキストを置換（layout 由来のスタイルを保持）。

    `ph.text = X` で全段落クリア + 新規テキスト挿入。layout / master の
    placeholder スタイル（フォント・サイズ・色）はそのまま継承される。
    """
    if ph is None:
        return
    ph.text = correct(text or "")


def _set_ph_text_styled(ph, text: str, *, size_pt=None, bold=None, color=None, font=None):
    """placeholder のテキストを置換 + 必要に応じてスタイル上書き。"""
    if ph is None:
        return
    _set_ph_text(ph, text)
    # 上書きスタイル
    tf = ph.text_frame
    if tf.paragraphs and tf.paragraphs[0].runs:
        run = tf.paragraphs[0].runs[0]
        if size_pt is not None:
            run.font.size = Pt(size_pt)
        if bold is not None:
            run.font.bold = bold
        if color is not None:
            run.font.color.rgb = color
        if font is not None:
            run.font.name = font


def _set_speaker_notes(slide, notes: list[str]) -> None:
    if not notes:
        return
    text = "\n".join(correct(n) for n in notes)
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


# ── スライドビルダー ───────────────────────────────────────


def build_title_slide(prs: Presentation, deck: Deck, layout) -> None:
    slide = prs.slides.add_slide(layout)
    # Wave Landscape の placeholder: Title (idx=0), SubTitle (idx=10), Date (idx=11)
    title_ph = _ph_by_name(slide, "Title_01", "Res_Title", "Title")
    sub_ph = _ph_by_name(slide, "SubTitle")
    date_ph = _ph_by_name(slide, "Date")

    _set_ph_text(title_ph, deck.title)
    sub_text = deck.subtitle or (
        "外部向けプレゼンテーション" if deck.audience.startswith("ext") else "社内資料"
    )
    _set_ph_text(sub_ph, sub_text)

    from datetime import date
    author_part = f"  /  {deck.author}" if deck.author else ""
    _set_ph_text(date_ph, f"{date.today().strftime('%Y.%m.%d')}{author_part}")


def build_agenda_slide(prs: Presentation, deck: Deck, layout, items: list[str]) -> None:
    slide = prs.slides.add_slide(layout)
    # アジェンダ layout は Res_Title_01 + Res_Chapter_0N + Res_Num_0N の placeholder セット
    # フォールバック: Title + Subtitle + Paragraph 系の場合は箇条書きで埋める
    title_ph = _ph_by_name(slide, "Title_01", "Res_Title", "Title")
    chapter_phs = _ph_all_by_name(slide, "Chapter_")
    num_phs = _ph_all_by_name(slide, "Num_")

    if chapter_phs:
        # 専用アジェンダ layout
        _set_ph_text(title_ph, "アジェンダ")
        for i, item in enumerate(items):
            if i >= len(chapter_phs):
                break
            _set_ph_text(chapter_phs[i], item)
            if i < len(num_phs):
                _set_ph_text(num_phs[i], str(i + 1))
        # 余りの placeholder は完全削除（残すとデフォルトテキストが出る）
        for j in range(len(items), len(chapter_phs)):
            _hide_shape(chapter_phs[j])
        for j in range(len(items), len(num_phs)):
            _hide_shape(num_phs[j])
    else:
        # フォールバック: Title + Subtitle + Paragraph 系
        _set_ph_text(title_ph, "アジェンダ")
        sub_ph = _ph_by_name(slide, "SubTitle")
        para_ph = _ph_by_name(slide, "Paraph_01", "Paraph")
        _set_ph_text(sub_ph, "本日お話しする内容")
        if para_ph is not None:
            tf = para_ph.text_frame
            tf.text = ""
            for i, item in enumerate(items):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.space_after = Pt(10)
                run_n = p.add_run()
                run_n.text = f"{i+1:02}.  "
                run_n.font.bold = True
                run_n.font.color.rgb = GOLD
                run_b = p.add_run()
                run_b.text = correct(item)


def _split_bullet(bullet: str) -> tuple[str, str]:
    """「課題 → 打ち手」「Label: 詳細」等を head / body に分割。"""
    for sep in (" → ", "→", " ⇒ ", "⇒"):
        if sep in bullet:
            head, _, body = bullet.partition(sep)
            return head.strip(), body.strip()
    for sep in ("：", ":"):
        if sep in bullet:
            head, _, body = bullet.partition(sep)
            if len(head.strip()) <= 30:  # 長すぎる head は分割しない
                return head.strip(), body.strip()
    return bullet.strip(), ""


def _detect_six_box_cells(slide):
    """Six Text Boxes layout から、(number_ph, header_ph, body_ph) の 6 ペアを返す。

    位置情報で動的にマッピングするためテンプレ依存しない。
    """
    placeholders = [
        ph for ph in slide.placeholders
        if ph.placeholder_format.idx not in (29, 30, 31)  # title/subtitle 除外
    ]
    # 番号 placeholder（width が小さい）/ header / body を分離
    numbers = [ph for ph in placeholders if (ph.width or 0) < 686000]  # < 0.75 inch
    not_numbers = [ph for ph in placeholders if (ph.width or 0) >= 686000]
    # 番号位置を基準に header（同じ top, 大きい left）と body（同じ left, 下）を見つける
    cells = []
    for num_ph in sorted(numbers, key=lambda p: (p.top or 0, p.left or 0)):
        num_top = num_ph.top or 0
        # header: 同じ top（±0.1 inch）かつ left が num_ph の右
        header_ph = None
        for ph in not_numbers:
            if abs((ph.top or 0) - num_top) < 91440 and (ph.left or 0) > (num_ph.left or 0):
                if header_ph is None or (ph.left or 0) < (header_ph.left or 0):
                    header_ph = ph
        if header_ph is None:
            continue
        # body: header の左と同じ、top が下
        body_ph = None
        for ph in not_numbers:
            if ph is header_ph:
                continue
            if abs((ph.left or 0) - (header_ph.left or 0)) < 91440 and (ph.top or 0) > (header_ph.top or 0):
                if body_ph is None or (ph.top or 0) < (body_ph.top or 0):
                    body_ph = ph
        cells.append((num_ph, header_ph, body_ph))
    # ソート: 行優先（top） → 列（left）
    cells.sort(key=lambda c: ((c[1].top or 0), (c[1].left or 0)))
    return cells


def build_content_six_text_boxes(prs: Presentation, deck: Deck, layout, section: Section) -> None:
    """6 セルグリッドで bullets をカード化（視認性高）。"""
    slide = prs.slides.add_slide(layout)
    title_ph = _ph_by_name(slide, "Title_01", "Res_Title", "Title")
    sub_ph = _ph_by_name(slide, "SubTitle")
    _set_ph_text(title_ph, section.title)
    if sub_ph is not None:
        _hide_shape(sub_ph)

    # Six Text Boxes は body に Res_Paraph_01 (idx=10) を含む。これは Title 配下では無く 1 セル目の body
    # なので _detect_six_box_cells で番号-ヘッダ-本文を位置で抽出する。
    cells = _detect_six_box_cells(slide)
    bullets = section.bullets

    for i, (num_ph, header_ph, body_ph) in enumerate(cells):
        if i < len(bullets):
            head, body = _split_bullet(bullets[i])
            _set_ph_text(num_ph, str(i + 1))
            _set_ph_text(header_ph, head)
            if body_ph is not None:
                _set_ph_text(body_ph, body)
        else:
            # 余ったセルは完全削除
            _hide_shape(num_ph)
            _hide_shape(header_ph)
            if body_ph is not None:
                _hide_shape(body_ph)

    _set_speaker_notes(slide, section.notes)


def build_content_two_paragraphs(prs: Presentation, deck: Deck, layout, section: Section) -> None:
    """2 列対比レイアウト。"""
    slide = prs.slides.add_slide(layout)
    title_ph = _ph_by_name(slide, "Title_01", "Res_Title", "Title")
    sub_ph = _ph_by_name(slide, "SubTitle")
    paras = _ph_all_by_name(slide, "Paraph_")

    _set_ph_text(title_ph, section.title)
    if sub_ph is not None:
        _hide_shape(sub_ph)

    bullets = section.bullets
    for i, ph in enumerate(paras):
        if i < len(bullets):
            head, body = _split_bullet(bullets[i])
            text = head if not body else f"{head}\n\n{body}"
            _set_ph_text(ph, text)
        else:
            _hide_shape(ph)

    _set_speaker_notes(slide, section.notes)


def build_content_paragraph(prs: Presentation, deck: Deck, layout, section: Section) -> None:
    """Title + Paragraph layout で番号付き強化。"""
    slide = prs.slides.add_slide(layout)
    title_ph = _ph_by_name(slide, "Title_01", "Res_Title", "Title")
    sub_ph = _ph_by_name(slide, "SubTitle")
    para_ph = _ph_by_name(slide, "Paraph_01", "Paraph")

    _set_ph_text(title_ph, section.title)
    if sub_ph is not None:
        _hide_shape(sub_ph)

    if para_ph is not None:
        tf = para_ph.text_frame
        tf.word_wrap = True
        tf.text = ""
        bullets = section.bullets if section.bullets else ["（内容を追加してください）"]
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(14)
            # 大きな Gold 番号（視覚アンカー）
            run_n = p.add_run()
            run_n.text = f"{i+1:02}    "
            run_n.font.bold = True
            run_n.font.size = Pt(20)
            run_n.font.color.rgb = GOLD
            run_n.font.name = "Inter"
            # 本文
            run_b = p.add_run()
            run_b.text = correct(b)
            run_b.font.size = Pt(16)
            run_b.font.name = "Noto Sans JP"

    _set_speaker_notes(slide, section.notes)


def build_content_slide(prs: Presentation, deck: Deck, layouts: dict, section: Section) -> None:
    """セクション内容に応じて最適な layout で本文スライドを構築。"""
    layout, kind = _pick_content_layout(layouts, section)
    if kind == "six":
        build_content_six_text_boxes(prs, deck, layout, section)
    elif kind == "two":
        build_content_two_paragraphs(prs, deck, layout, section)
    else:
        build_content_paragraph(prs, deck, layout, section)


def build_kpi_slide(prs: Presentation, deck: Deck, layout, section: Section) -> None:
    slide = prs.slides.add_slide(layout)
    title_ph = _ph_by_name(slide, "Title_01", "Res_Title", "Title")
    sub_ph = _ph_by_name(slide, "SubTitle")
    _set_ph_text(title_ph, section.title or "KPI ハイライト")
    if section.bullets:
        _set_ph_text(sub_ph, section.bullets[0])
    else:
        _set_ph_text(sub_ph, "重点指標と達成度")

    # Six Text Boxes layout の場合: 6 つの番号 + ラベル + 値 placeholder ペア
    # placeholder 名: 'Text Placeholder N' / 'Res_Paraph_01' なども
    # 単純化: layout に Text Placeholder N が複数ある場合、KPI 数 × 2 ペアで埋める
    text_phs = [
        ph for ph in slide.placeholders
        if "Text Placeholder" in (ph.name or "") or "Espace réservé du texte" in (ph.name or "")
    ]
    # idx 順
    text_phs.sort(key=lambda p: p.placeholder_format.idx)

    # KPI を順に埋める: 番号 placeholder と Label/Value placeholder にマッピング
    if len(text_phs) >= 2 and section.kpis:
        # 最初の N 個を「番号」placeholder（小さめサイズ）として、その後を value+label として使う
        # 簡易戦略: テキスト placeholder の前半を「label/value」対応に。
        # Six Text Boxes は 6 ボックス × (ラベル+本文) = 12 placeholder あるので半分使う
        pairs_count = min(len(section.kpis), max(1, len(text_phs) // 3))
        ph_iter = iter(text_phs)
        for i, (label, value) in enumerate(section.kpis):
            try:
                title_box = next(ph_iter)
                _set_ph_text_styled(title_box, label, size_pt=14, bold=True, color=NAVY)
                body_box = next(ph_iter)
                _set_ph_text_styled(body_box, value, size_pt=44, bold=True, color=GOLD, font="Inter")
            except StopIteration:
                break
        # 残りは空文字
        for ph in ph_iter:
            _set_ph_text(ph, "")
    elif section.kpis:
        # フォールバック: Centered Text Box 2 や Title+Paragraph
        para_ph = _ph_by_name(slide, "Paraph_01", "Paraph")
        if para_ph is not None:
            tf = para_ph.text_frame
            tf.word_wrap = True
            tf.text = ""
            for i, (label, value) in enumerate(section.kpis):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.space_after = Pt(8)
                run_v = p.add_run()
                run_v.text = f"{value}  "
                run_v.font.bold = True
                run_v.font.size = Pt(36)
                run_v.font.color.rgb = GOLD
                run_v.font.name = "Inter"
                run_l = p.add_run()
                run_l.text = label
                run_l.font.size = Pt(16)
                run_l.font.color.rgb = NAVY

    _set_speaker_notes(slide, section.notes)


def build_closing_slide(prs: Presentation, deck: Deck, layout, section: Section) -> None:
    slide = prs.slides.add_slide(layout)
    # 'Thank you' / 'End with Voice' は placeholder が無いケースが多い
    # placeholder があれば埋める、無ければそのまま（layout のデザインだけで完結）
    title_ph = _ph_by_name(slide, "Title_01", "Res_Title", "Title")
    sub_ph = _ph_by_name(slide, "SubTitle")
    if title_ph is not None:
        _set_ph_text(title_ph, section.title or "Thank you")
    if sub_ph is not None:
        if section.bullets:
            _set_ph_text(sub_ph, " · ".join(correct(b) for b in section.bullets))
        else:
            _set_ph_text(sub_ph, "AKKODiS — Engineering future, together.")
    _set_speaker_notes(slide, section.notes)


# ── ビルド本体 ──────────────────────────────────────────────


def build(deck: Deck, template: Path) -> Presentation:
    prs = Presentation(str(template))
    xml_slides = prs.slides._sldIdLst
    slides_part = prs.part
    for sld in list(xml_slides):
        rId = sld.get(qn("r:id"))
        if rId:
            slides_part.drop_rel(rId)
        xml_slides.remove(sld)

    # KPI 数を見て KPI レイアウトを決める
    kpi_max = max(
        (len(s.kpis) for s in deck.sections if s.kind == "kpi"), default=0
    )
    layouts = _select_layouts(prs, n_kpi=kpi_max)

    has_agenda = any(s.kind == "agenda" for s in deck.sections)
    insert_agenda = (not has_agenda) and len(
        [s for s in deck.sections if s.kind != "agenda"]
    ) >= 2

    build_title_slide(prs, deck, layouts["title"])

    if insert_agenda:
        agenda_items = [s.title for s in deck.sections if s.kind != "agenda"]
        build_agenda_slide(prs, deck, layouts["agenda"] or layouts["content_paragraph"], agenda_items)

    for sec in deck.sections:
        if sec.kind == "agenda":
            items = sec.bullets if sec.bullets else [
                s.title for s in deck.sections if s != sec and s.kind != "agenda"
            ]
            build_agenda_slide(prs, deck, layouts["agenda"] or layouts["content_paragraph"], items)
        elif sec.kind == "kpi":
            build_kpi_slide(prs, deck, layouts["kpi"] or layouts["content_paragraph"], sec)
        elif sec.kind == "closing":
            build_closing_slide(prs, deck, layouts["closing"] or layouts["title"], sec)
        else:
            build_content_slide(prs, deck, layouts, sec)

    return prs


def resolve_template(template: Path | None, deck: Deck) -> Path:
    if template is not None:
        return template
    here = Path(__file__).resolve().parent.parent
    aud = "external" if deck.audience.startswith("ext") else "internal"
    tone = "dark" if deck.tone == "dark" else "white"
    return here / "templates" / f"{aud}-{tone}.pptx"


def main(argv: list[str]) -> int:
    ap = argparse.ArgumentParser(description="AKKODiS ブランド準拠 PPTX 生成")
    ap.add_argument("--template", type=Path, default=None)
    ap.add_argument("--input", type=Path, required=True)
    ap.add_argument("--output", type=Path, required=True)
    args = ap.parse_args(argv)

    text = args.input.read_text(encoding="utf-8")
    deck = parse_input(text)
    template = resolve_template(args.template, deck)
    if not template.exists():
        print(f"[error] テンプレートが見つかりません: {template}", file=sys.stderr)
        return 1
    prs = build(deck, template)
    args.output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(args.output))
    print(f"[ok] {args.output}  ({len(prs.slides)} slides)")
    return 0


if __name__ == "__main__":
    sys.exit(main(sys.argv[1:]))
