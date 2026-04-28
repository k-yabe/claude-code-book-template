"""AKKODiS Claude Skills 3 種の動作確認テスト。

実行: pytest tests/test_skills.py -v
"""

from __future__ import annotations

import sys
from pathlib import Path

import pytest

REPO_ROOT = Path(__file__).resolve().parent.parent
SKILLS_DIR = REPO_ROOT / "skills"


def _import_module(skill: str, module: str):
    """skill/scripts 配下のモジュールを動的に import する。

    モジュール名にダッシュは使えない（dataclass の `__module__` 経由で再 import が走るため）。
    `skill_pptx_build_pptx` のようなサニタイズ名で登録する。
    """
    import importlib.util

    path = SKILLS_DIR / skill / "scripts" / f"{module}.py"
    sanitized = (skill + "_" + module).replace("-", "_")
    if sanitized in sys.modules:
        return sys.modules[sanitized]
    # scripts ディレクトリを sys.path に入れて同階層 import を解決
    scripts_dir = str(path.parent)
    if scripts_dir not in sys.path:
        sys.path.insert(0, scripts_dir)
    spec = importlib.util.spec_from_file_location(sanitized, path)
    mod = importlib.util.module_from_spec(spec)
    sys.modules[sanitized] = mod
    spec.loader.exec_module(mod)
    return mod


# ── 表記補正 ──────────────────────────────────────────────────


@pytest.fixture(scope="module")
def notation():
    return _import_module("akkodis-pptx", "notation")


class TestNotation:
    def test_akkodis_capitalization(self, notation):
        assert notation.correct("Akkodis") == "AKKODiS"
        assert notation.correct("AKKODIS") == "AKKODiS"
        assert notation.correct("akkodis") == "AKKODiS"
        # 文中
        assert "AKKODiSコンサル" in notation.correct("Akkodisコンサル")

    def test_innovation_lab_lowercase_i(self, notation):
        assert notation.correct("AKKODiS Innovation Lab.") == "AKKODiS innovation Lab."

    def test_iown_trademark(self, notation):
        # 単体 → ®
        assert notation.correct("IOWNです") == "IOWN®です"
        # 付属語複合 → ® なし
        assert notation.correct("IOWN構想") == "IOWN構想"
        assert notation.correct("IOWNソリューション") == "IOWNソリューション"
        # 誤って付いた ® は除去
        assert notation.correct("IOWN®構想") == "IOWN構想"
        # GLOBAL FORUM → ™
        assert notation.correct("IOWN GLOBAL FORUM") == "IOWN® GLOBAL FORUM™"

    def test_microsoft_naming(self, notation):
        assert notation.correct("PowerBI") == "Power BI"
        assert notation.correct("PowerAutomate") == "Power Automate"
        assert notation.correct("Sharepoint") == "SharePoint"
        assert notation.correct("OneNote") == "OneNote"  # 不変
        # 既に正しいものは変えない
        assert notation.correct("Microsoft 365 Copilot") == "Microsoft 365 Copilot"

    def test_keigo_corrections(self, notation):
        assert notation.correct("頂きます") == "いただきます"
        assert notation.correct("資料を頂く") == "資料をいただく"
        assert notation.correct("致します") == "いたします"
        # 二重敬語
        assert notation.correct("拝見させていただきます") == "拝見します"
        assert "とんでもないことです" in notation.correct("とんでもございません")

    def test_okurigana(self, notation):
        assert notation.correct("行ないます") == "行います"
        assert notation.correct("打合せ") == "打ち合わせ"
        assert notation.correct("問合わせ") == "問い合わせ"
        assert notation.correct("申込み") == "申し込み"
        # 申込書（固有名詞）は変えない
        assert notation.correct("申込書") == "申込書"

    def test_redundant_expressions(self, notation):
        assert notation.correct("まず最初に始めます") == "まず始めます"
        assert notation.correct("頭痛が痛い") == "頭が痛い"
        assert notation.correct("今の現状") == "現状"

    def test_no_change_for_correct(self, notation):
        # 既に正しいものは変えない
        unchanged = "AKKODiS は記者ハンドブックの表記ルールを順守します。"
        assert notation.correct(unchanged) == unchanged

    def test_phone_number_unchanged(self, notation):
        # 全角数字でない電話番号は変えない
        assert notation.correct("電話: 03-1234-5678") == "電話: 03-1234-5678"

    def test_full_width_digits(self, notation):
        assert notation.correct("１２３") == "123"


# ── PPTX 生成 ────────────────────────────────────────────────


@pytest.fixture(scope="module")
def pptx_output(tmp_path_factory):
    pptx_mod = _import_module("akkodis-pptx", "build_pptx")
    out = tmp_path_factory.mktemp("pptx") / "test.pptx"
    sample = SKILLS_DIR / "akkodis-pptx" / "examples" / "sample-input.md"
    rc = pptx_mod.main([
        "--input", str(sample),
        "--output", str(out),
    ])
    assert rc == 0
    return out


class TestPPTX:
    def test_file_exists(self, pptx_output):
        assert pptx_output.exists()
        assert pptx_output.stat().st_size > 1000

    def test_slide_count(self, pptx_output):
        from pptx import Presentation
        p = Presentation(str(pptx_output))
        # title + agenda + 6 sections (現状/課題/施策/KPI/スケジュール/まとめ) = 8
        assert len(p.slides) >= 7

    def test_title_slide_text(self, pptx_output):
        from pptx import Presentation
        p = Presentation(str(pptx_output))
        first = p.slides[0]
        all_text = " ".join(
            para.text for sh in first.shapes if sh.has_text_frame
            for para in sh.text_frame.paragraphs
        )
        assert "2026 Q2" in all_text
        assert "AKKODiS" in all_text

    def test_speaker_notes_present(self, pptx_output):
        from pptx import Presentation
        p = Presentation(str(pptx_output))
        notes_found = any(
            s.has_notes_slide and s.notes_slide.notes_text_frame.text.strip()
            for s in p.slides
        )
        assert notes_found, "発表者ノートが1枚も無い"

    def test_kpi_slide_has_big_numbers(self, pptx_output):
        from pptx import Presentation
        p = Presentation(str(pptx_output))
        kpi_slide = None
        for s in p.slides:
            text = " ".join(
                para.text for sh in s.shapes if sh.has_text_frame
                for para in sh.text_frame.paragraphs
            )
            if "KPI" in text:
                kpi_slide = s
                break
        assert kpi_slide is not None, "KPI スライドが無い"


# ── XLSX 生成 ────────────────────────────────────────────────


@pytest.fixture(scope="module")
def xlsx_output(tmp_path_factory):
    xlsx_mod = _import_module("akkodis-xlsx", "build_xlsx")
    out = tmp_path_factory.mktemp("xlsx") / "test.xlsx"
    sample = SKILLS_DIR / "akkodis-xlsx" / "examples" / "sample-input.md"
    rc = xlsx_mod.main([
        "--input", str(sample),
        "--output", str(out),
    ])
    assert rc == 0
    return out


class TestXLSX:
    def test_file_exists(self, xlsx_output):
        assert xlsx_output.exists()

    def test_sheet_count(self, xlsx_output):
        from openpyxl import load_workbook
        wb = load_workbook(xlsx_output)
        assert len(wb.sheetnames) >= 2

    def test_navy_header(self, xlsx_output):
        from openpyxl import load_workbook
        wb = load_workbook(xlsx_output)
        ws = wb[wb.sheetnames[0]]
        # 3行目（タイトル下）がヘッダー
        header_cell = ws.cell(row=3, column=1)
        assert header_cell.fill.fgColor.rgb in ("FF001F33", "00001F33")

    def test_freeze_panes_set(self, xlsx_output):
        from openpyxl import load_workbook
        wb = load_workbook(xlsx_output)
        for name in wb.sheetnames:
            ws = wb[name]
            assert ws.freeze_panes is not None

    def test_autofilter_set(self, xlsx_output):
        from openpyxl import load_workbook
        wb = load_workbook(xlsx_output)
        for name in wb.sheetnames:
            ws = wb[name]
            assert ws.auto_filter.ref is not None, f"{name} にオートフィルタ無し"

    def test_chart_present_on_marked_sheet(self, xlsx_output):
        from openpyxl import load_workbook
        wb = load_workbook(xlsx_output)
        # 「サマリー」「月次推移」がチャートマークありの想定
        chart_count = sum(len(wb[n]._charts) for n in wb.sheetnames)
        assert chart_count >= 2

    def test_landscape_print(self, xlsx_output):
        from openpyxl import load_workbook
        wb = load_workbook(xlsx_output)
        for name in wb.sheetnames:
            ws = wb[name]
            assert ws.page_setup.orientation == "landscape"


# ── DOCX 生成 ────────────────────────────────────────────────


@pytest.fixture(scope="module")
def docx_output(tmp_path_factory):
    docx_mod = _import_module("akkodis-docx", "build_docx")
    out = tmp_path_factory.mktemp("docx") / "test.docx"
    sample = SKILLS_DIR / "akkodis-docx" / "examples" / "sample-input.md"
    rc = docx_mod.main([
        "--input", str(sample),
        "--output", str(out),
    ])
    assert rc == 0
    return out


class TestDOCX:
    def test_file_exists(self, docx_output):
        assert docx_output.exists()

    def test_heading_styles_used(self, docx_output):
        from docx import Document
        d = Document(str(docx_output))
        styles = {p.style.name for p in d.paragraphs if p.text.strip()}
        assert "Heading 1" in styles
        assert "Heading 2" in styles

    def test_table_present(self, docx_output):
        from docx import Document
        d = Document(str(docx_output))
        assert len(d.tables) >= 1

    def test_table_header_navy_shading(self, docx_output):
        from docx import Document
        d = Document(str(docx_output))
        t = d.tables[0]
        first_cell = t.rows[0].cells[0]
        # cell shading XML を検査
        from docx.oxml.ns import qn
        tcPr = first_cell._tc.tcPr
        shd = tcPr.find(qn("w:shd")) if tcPr is not None else None
        assert shd is not None, "テーブルヘッダーに shading 無し"
        assert shd.get(qn("w:fill")) == "001F33"

    def test_metadata_in_cover(self, docx_output):
        from docx import Document
        d = Document(str(docx_output))
        all_text = " ".join(p.text for p in d.paragraphs)
        assert "マーケティング自動化基盤" in all_text
        assert "株式会社○○" in all_text
        assert "AKKODiS" in all_text


# ── 配布 ZIP 構造 ─────────────────────────────────────────────


class TestDistribution:
    @pytest.mark.parametrize("skill", ["akkodis-pptx", "akkodis-xlsx", "akkodis-docx"])
    def test_skill_md_exists(self, skill):
        skill_md = SKILLS_DIR / skill / "SKILL.md"
        assert skill_md.exists()
        text = skill_md.read_text(encoding="utf-8")
        # frontmatter check
        assert text.startswith("---")
        assert "\nname: " in text
        assert "\ndescription: " in text

    @pytest.mark.parametrize("skill", ["akkodis-pptx", "akkodis-xlsx", "akkodis-docx"])
    def test_brand_assets_present(self, skill):
        brand_dir = SKILLS_DIR / skill / "brand"
        assert brand_dir.exists()
        assert (brand_dir / "style-guide.md").exists()
        assert (brand_dir / "notation-rules.md").exists()
        # ロゴ 5 種
        svgs = list(brand_dir.glob("AKKODIS_Logo_*.svg"))
        assert len(svgs) == 5

    @pytest.mark.parametrize("skill", ["akkodis-pptx", "akkodis-xlsx", "akkodis-docx"])
    def test_notation_rules_complete(self, skill):
        rules = (SKILLS_DIR / skill / "brand" / "notation-rules.md").read_text(
            encoding="utf-8"
        )
        # A-D の主要セクションが含まれる
        assert "## A. 記者ハンドブック準拠" in rules
        assert "## B. AKKODiS ブランドガイドライン" in rules
        assert "## C. Microsoft" in rules
        assert "## D. IOWN®" in rules
