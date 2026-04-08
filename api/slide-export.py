"""
Vercel Python Serverless Function: PPTX生成（python-pptx add_slide方式）
テンプレートのslide_layoutを使い、テキスト + ネイティブチャート + テーブル + フロー図形 + 画像を生成。
"""
from http.server import BaseHTTPRequestHandler
import json
import os
import io
import re
import base64
import urllib.request
import urllib.error
from concurrent.futures import ThreadPoolExecutor, as_completed
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData

TEMPLATE_DIR = os.path.join(os.path.dirname(os.path.dirname(__file__)), 'apps', 'slide-maker', 'templates')

UNSPLASH_ACCESS_KEY = os.environ.get('UNSPLASH_ACCESS_KEY', '')

LAYOUT_NAMES = {
    'cover': 'Wave in the corners',
    'chapter': 'Chapter - Mesh Gold',
    'agenda': 'Agenda - Computer',
    'content': 'Title, Subtitle and one Paragrah',
    'two-column': 'Two Paragraphs with Blue Line',
    'centered-text': 'Centered Text Box',
    'text-left-picture-right': 'Text left - Picture Right',
    'text-right-picture-left': 'Text Right- Picture Left',
    'content-left-full-picture-right': 'Content Left - Full Picture Right',
    'large-image-right': 'Large Image - Right',
    'picture-fullscreen': 'Picture full screen with Title block',
    'content-with-chart': 'Title, Subtitle and one Paragrah',
    'content-with-flow': 'Title, Subtitle and one Paragrah',
    'sixbox': 'Six Text Boxes',
    'comparison': 'Title, Subtitle and one Paragrah',
    'quote': '1_Quote Content with animation',
    'closing': 'Thank you',
}

# 画像系レイアウト — Unsplash未設定時はcontentにフォールバック
IMAGE_LAYOUTS = {
    'text-left-picture-right', 'text-right-picture-left',
    'content-left-full-picture-right', 'large-image-right', 'picture-fullscreen',
}

NAVY = RGBColor(0x00, 0x1F, 0x33)
GOLD = RGBColor(0xFF, 0xB8, 0x1C)
CYAN = RGBColor(0x00, 0xBF, 0xD6)
GRAY = RGBColor(0x8A, 0x9B, 0xB0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CHART_COLORS = [NAVY, GOLD, CYAN, GRAY, RGBColor(0x4A, 0x55, 0x68), RGBColor(0xE6, 0x7E, 0x22)]


def fetch_unsplash_image(query, timeout=5):
    """Unsplash APIから画像を取得してBytesIOで返す。失敗時はNone。"""
    if not UNSPLASH_ACCESS_KEY or not query:
        return None
    try:
        search_url = f'https://api.unsplash.com/photos/random?query={urllib.request.quote(query)}&orientation=landscape&w=1200'
        req = urllib.request.Request(search_url, headers={
            'Authorization': f'Client-ID {UNSPLASH_ACCESS_KEY}',
            'Accept-Version': 'v1',
        })
        with urllib.request.urlopen(req, timeout=timeout) as resp:
            data = json.loads(resp.read())
            img_url = data.get('urls', {}).get('regular', '')
            if not img_url:
                return None
        img_req = urllib.request.Request(img_url)
        with urllib.request.urlopen(img_req, timeout=timeout) as img_resp:
            return io.BytesIO(img_resp.read())
    except Exception:
        return None


def fetch_images_batch(slides_data):
    """複数スライドの画像を並列取得。{index: BytesIO} を返す。"""
    image_tasks = {}
    for i, sd in enumerate(slides_data):
        if sd.get('layout') in IMAGE_LAYOUTS and sd.get('imageQuery'):
            image_tasks[i] = sd['imageQuery']

    if not image_tasks:
        return {}

    results = {}
    with ThreadPoolExecutor(max_workers=5) as executor:
        futures = {executor.submit(fetch_unsplash_image, q): idx for idx, q in image_tasks.items()}
        for future in as_completed(futures):
            idx = futures[future]
            try:
                img = future.result()
                if img:
                    results[idx] = img
            except Exception:
                pass
    return results


def insert_picture_to_placeholder(slide, ph_idx, image_stream):
    """PICTURE placeholderに画像を挿入。"""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == ph_idx:
            try:
                ph.insert_picture(image_stream)
                return True
            except Exception:
                pass
    return False


def _set(ph_map, idx, text, font_size_pt=None):
    if idx in ph_map:
        try:
            # python-pptxは自動でXMLエスケープするのでstr変換のみ
            ph_map[idx].text = str(text).strip() if text else ''
            if font_size_pt is not None and ph_map[idx].text:
                for para in ph_map[idx].text_frame.paragraphs:
                    for run in para.runs:
                        run.font.size = _auto_font_size(text, font_size_pt)
        except Exception:
            pass


JP_MULTIPLIERS = {'万': 10000, '億': 100000000, '兆': 1000000000000}


def _parse_numeric(v):
    """数値部分を抽出。「1.2万円」→12000、「$500」→500、「100人」→100。"""
    if isinstance(v, (int, float)):
        return float(v)
    s = str(v).strip()
    if not s:
        return 0
    # 日本語単位の倍率を検出
    multiplier = 1
    for unit_char, mult in JP_MULTIPLIERS.items():
        if unit_char in s:
            multiplier = mult
            s = s.replace(unit_char, '')
            break
    # 数値部分のみ抽出（マイナス・小数点を許容）
    cleaned = re.sub(r'[^\d.\-]', '', s)
    if not cleaned or cleaned in ('.', '-', '-.'):
        return 0
    try:
        return float(cleaned) * multiplier
    except (ValueError, TypeError):
        return 0


def _auto_font_size(text, base_pt, long_threshold=40, min_pt=8):
    """テキスト長に応じてフォントサイズを自動調整する。"""
    length = len(str(text)) if text else 0
    if length <= long_threshold:
        return Pt(base_pt)
    # 長いテキストは段階的に縮小
    reduction = min((length - long_threshold) // 20, (base_pt - min_pt) // 2)
    return Pt(max(min_pt, base_pt - reduction * 2))


def add_chart(slide, chart_data):
    """スライドにネイティブPPTXチャートを追加（高品質スタイリング）"""
    ctype = chart_data.get('type', 'bar')
    chart_type_map = {
        'bar': XL_CHART_TYPE.COLUMN_CLUSTERED,
        'line': XL_CHART_TYPE.LINE_MARKERS,
        'pie': XL_CHART_TYPE.PIE,
    }
    xl_type = chart_type_map.get(ctype, XL_CHART_TYPE.COLUMN_CLUSTERED)

    cd = CategoryChartData()
    cd.categories = chart_data.get('labels', [])
    values = chart_data.get('data', [])
    num_values = []
    for v in values:
        num_values.append(_parse_numeric(v))
    cd.add_series(chart_data.get('title', ''), num_values)

    x, y = Inches(0.8), Inches(2.6)
    cx, cy = Inches(8.4), Inches(4.2)
    chart_frame = slide.shapes.add_chart(xl_type, x, y, cx, cy, cd)
    chart = chart_frame.chart

    # 共通スタイリング
    chart.has_legend = (ctype == 'pie')
    chart.font.size = Pt(10)
    chart.font.color.rgb = NAVY

    if chart.series:
        series = chart.series[0]
        if ctype == 'bar':
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = NAVY
            # データラベル表示
            series.has_data_labels = True
            series.data_labels.font.size = Pt(9)
            series.data_labels.font.color.rgb = NAVY
            series.data_labels.font.bold = True
            series.data_labels.number_format = '#,##0'
        elif ctype == 'line':
            series.format.line.color.rgb = NAVY
            series.format.line.width = Pt(2.5)
            series.smooth = True
            series.has_data_labels = True
            series.data_labels.font.size = Pt(8)
            series.data_labels.font.color.rgb = NAVY
        elif ctype == 'pie':
            for i, point in enumerate(series.points):
                point.format.fill.solid()
                point.format.fill.fore_color.rgb = CHART_COLORS[i % len(CHART_COLORS)]
            series.has_data_labels = True
            series.data_labels.font.size = Pt(9)
            series.data_labels.show_percentage = True
            series.data_labels.show_category_name = True

    # 軸スタイリング（pie以外）
    if ctype != 'pie':
        if hasattr(chart, 'value_axis'):
            chart.value_axis.has_major_gridlines = True
            chart.value_axis.major_gridlines.format.line.color.rgb = RGBColor(0xE0, 0xE4, 0xE8)
            chart.value_axis.format.line.color.rgb = RGBColor(0xE0, 0xE4, 0xE8)
        if hasattr(chart, 'category_axis'):
            chart.category_axis.format.line.color.rgb = RGBColor(0xE0, 0xE4, 0xE8)

    # 単位ラベル
    unit = chart_data.get('unit', '')
    if unit and ctype != 'pie' and hasattr(chart, 'value_axis'):
        try:
            chart.value_axis.axis_title.text_frame.text = unit
            chart.value_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(9)
            chart.value_axis.axis_title.text_frame.paragraphs[0].font.color.rgb = GRAY
        except Exception:
            pass


def add_table(slide, table_data):
    """ネイティブテーブル（ゼブラストライプ + 列幅自動調整 + 見やすいフォーマット）"""
    headers = table_data.get('headers', [])
    rows_data = table_data.get('rows', [])
    if not headers:
        return

    cols = len(headers)
    row_count = 1 + len(rows_data)
    row_h = Inches(0.55)
    x, y = Inches(0.8), Inches(2.3)
    cx = Inches(8.4)
    cy = row_h * row_count
    table_shape = slide.shapes.add_table(row_count, cols, x, y, cx, cy)
    table = table_shape.table

    # 列幅自動調整: 各列の最大文字数に基づいて比例配分
    col_max_lens = []
    for j in range(cols):
        max_len = len(str(headers[j]))
        for row in rows_data:
            if j < len(row):
                max_len = max(max_len, len(str(row[j])))
        col_max_lens.append(max(max_len, 1))  # 最低1
    total_len = sum(col_max_lens)
    for j in range(cols):
        table.columns[j].width = int(cx * col_max_lens[j] / total_len)

    LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)

    # ヘッダー行
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = str(h)
        for para in cell.text_frame.paragraphs:
            para.font.size = Pt(14)
            para.font.bold = True
            para.font.color.rgb = WHITE
            para.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        cell.margin_top = Pt(6)
        cell.margin_bottom = Pt(6)

    # データ行（ゼブラストライプ + カラムパディング）
    for i, row in enumerate(rows_data):
        # カラム数が足りない場合はパディング
        while len(row) < cols:
            row.append('')
        for j, val in enumerate(row):
            if j < cols:
                cell = table.cell(i + 1, j)
                cell.text = str(val)
                for para in cell.text_frame.paragraphs:
                    para.font.size = _auto_font_size(val, 11)
                    para.font.color.rgb = NAVY
                cell.margin_top = Pt(4)
                cell.margin_bottom = Pt(4)
                # 偶数行に背景色
                if i % 2 == 1:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = LIGHT_BG


def add_flow(slide, steps):
    """フロー図（ステップ番号 + ネイビーボックス + ゴールド矢印、ステップ数に応じて自動調整）"""
    if not steps:
        return
    n = len(steps)
    # ステップ数に応じたレスポンシブサイズ計算
    avail_w = Inches(8.4)
    gap_w = Inches(0.18) if n > 4 else Inches(0.22)
    box_w = int((avail_w - gap_w * (n - 1)) / n)
    max_box_w = Inches(1.8)
    min_box_w = Inches(0.7)
    box_w = max(min_box_w, min(max_box_w, box_w))
    total = box_w * n + gap_w * (n - 1)
    start_x = Inches(0.8) + (avail_w - total) // 2
    box_h = Inches(0.85)
    y = Inches(3.8)

    # ステップ数に応じたフォントサイズ
    if n <= 3:
        step_font = 11
        num_font = 11
    elif n <= 5:
        step_font = 9
        num_font = 10
    else:
        step_font = 8
        num_font = 9

    for i, step in enumerate(steps):
        x = start_x + i * (box_w + gap_w)

        # ステップ番号（上）
        num_shape = slide.shapes.add_shape(1, x + box_w // 2 - Inches(0.18), y - Inches(0.5), Inches(0.36), Inches(0.36))
        num_shape.fill.solid()
        num_shape.fill.fore_color.rgb = GOLD
        num_shape.line.fill.background()
        np = num_shape.text_frame.paragraphs[0]
        np.text = str(i + 1)
        np.font.size = Pt(num_font)
        np.font.color.rgb = NAVY
        np.font.bold = True
        np.alignment = PP_ALIGN.CENTER

        # メインボックス
        shape = slide.shapes.add_shape(1, x, y, box_w, box_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = NAVY
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = str(step)
        p.font.size = _auto_font_size(step, step_font, long_threshold=15, min_pt=7)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # 矢印（最後以外）
        if i < n - 1:
            ax = x + box_w
            arrow = slide.shapes.add_shape(13, ax + Inches(0.04), y + box_h // 3, gap_w - Inches(0.08), box_h // 3)
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = GOLD
            arrow.line.fill.background()


def add_feature_cards(slide, body_text, x, y, width):
    """箇条書きを色付きカードとして配置（参考PPTXスライド9のパターン）"""
    lines = [l.strip() for l in str(body_text).split('\n') if l.strip()]
    if not lines:
        return
    n = len(lines)
    card_h = Inches(0.75)
    gap = Inches(0.1)
    bar_w = Inches(0.06)

    for i, line in enumerate(lines):
        cy = y + i * (card_h + gap)
        text = line.lstrip('・-•● ')

        # 左のカラーバー（ゴールド）
        bar = slide.shapes.add_shape(1, x, cy, bar_w, card_h)
        bar.fill.solid()
        bar.fill.fore_color.rgb = GOLD
        bar.line.fill.background()

        # カード背景
        card_bg = slide.shapes.add_shape(1, x + bar_w, cy, width - bar_w, card_h)
        card_bg.fill.solid()
        card_bg.fill.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFA)
        card_bg.line.fill.background()

        # テキスト
        txbox = slide.shapes.add_textbox(x + bar_w + Inches(0.2), cy + Inches(0.1), width - bar_w - Inches(0.4), card_h - Inches(0.2))
        tf = txbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(10)
        p.font.color.rgb = NAVY
        p.font.bold = False


def add_section_dividers(slide, x, y, width, count):
    """水平区切り線を追加"""
    for i in range(count):
        ly = y + i * Inches(1.2)
        line = slide.shapes.add_shape(1, x, ly, width, Inches(0.015))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(0xE0, 0xE4, 0xE8)
        line.line.fill.background()


def add_accent_line(slide, x, y, width):
    """ゴールドのアクセントライン"""
    line = slide.shapes.add_shape(1, x, y, width, Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = GOLD
    line.line.fill.background()


def add_bullet_icons(slide, body_text, x, y, width):
    """箇条書きをアイコンバレット付きのshapeとして描画"""
    lines = [l.strip() for l in str(body_text).split('\n') if l.strip()]
    line_h = Inches(0.38)
    icon_size = Inches(0.22)

    for i, line in enumerate(lines):
        ly = y + i * line_h
        text = line.lstrip('・-•● ')

        # バレットアイコン（丸）
        icon = slide.shapes.add_shape(9, x, ly + Inches(0.04), icon_size, icon_size)  # 9 = oval
        icon.fill.solid()
        icon.fill.fore_color.rgb = GOLD
        icon.line.fill.background()

        # テキスト
        txbox = slide.shapes.add_textbox(x + Inches(0.35), ly, width - Inches(0.35), line_h)
        tf = txbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(11)
        p.font.color.rgb = NAVY


def add_metric_callout(slide, value, label, x, y, width=Inches(1.8), height=Inches(1.2)):
    """大きな数字 + ラベルのメトリクスカード"""
    # 背景ボックス
    box = slide.shapes.add_shape(1, x, y, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFA)
    box.line.fill.background()

    # 数字（大きく）
    num_box = slide.shapes.add_textbox(x, y + Inches(0.1), width, Inches(0.6))
    p = num_box.text_frame.paragraphs[0]
    p.text = str(value)
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.CENTER

    # ラベル（小さく）
    lbl_box = slide.shapes.add_textbox(x, y + Inches(0.7), width, Inches(0.4))
    p2 = lbl_box.text_frame.paragraphs[0]
    p2.text = str(label)
    p2.font.size = Pt(9)
    p2.font.color.rgb = GRAY
    p2.alignment = PP_ALIGN.CENTER

    # ゴールドトップライン
    add_accent_line(slide, x, y, width)


def set_rich_body(ph, body_text, base_size=14):
    """本文をリッチテキスト（太字/サイズ使い分け）で設定"""
    if not body_text:
        ph.text = ''
        return

    tf = ph.text_frame
    tf.clear()
    lines = str(body_text).split('\n')

    for i, line in enumerate(lines):
        line = line.strip()
        if not line:
            continue
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()

        # 見出し行（【】で囲まれている）
        if line.startswith('【') and '】' in line:
            run = p.add_run()
            run.text = line
            run.font.size = _auto_font_size(line, base_size + 2)
            run.font.bold = True
            run.font.color.rgb = NAVY
        # 箇条書き行
        elif line.startswith(('・', '- ', '• ', '● ')):
            run = p.add_run()
            run.text = line
            run.font.size = _auto_font_size(line, base_size)
            run.font.color.rgb = RGBColor(0x2D, 0x34, 0x36)
            p.space_before = Pt(2)
        # 通常行
        else:
            run = p.add_run()
            run.text = line
            run.font.size = _auto_font_size(line, base_size)
            run.font.color.rgb = RGBColor(0x2D, 0x34, 0x36)


def apply_data(slide, layout, data):
    ph_map = {ph.placeholder_format.idx: ph for ph in slide.placeholders}

    if layout == 'cover':
        _set(ph_map, 0, data.get('title', ''), font_size_pt=36)
        _set(ph_map, 10, data.get('subtitle', ''), font_size_pt=18)
        _set(ph_map, 11, data.get('date', ''))

    elif layout == 'chapter':
        _set(ph_map, 12, data.get('title', ''))
        _set(ph_map, 13, data.get('number', ''))

    elif layout == 'agenda':
        items = data.get('items', [])
        _set(ph_map, 24, data.get('title', 'Agenda'))
        for j in range(6):
            _set(ph_map, 12 + j * 2, items[j] if j < len(items) else '')
            _set(ph_map, 13 + j * 2, str(j + 1) if j < len(items) else '')

    elif layout == 'content':
        _set(ph_map, 29, data.get('title', ''), font_size_pt=28)
        _set(ph_map, 30, '')
        body = data.get('body', '')
        lines = [l.strip() for l in str(body).split('\n') if l.strip()] if body else []
        has_bullets = any(l.startswith(('・', '- ', '• ', '● ')) for l in lines)

        if has_bullets and len(lines) >= 3:
            # 3行以上の箇条書き → feature cards（参考PPTXパターン）
            _set(ph_map, 32, '')
            add_feature_cards(slide, body, Inches(0.8), Inches(2.1), Inches(8.4))
            add_accent_line(slide, Inches(0.8), Inches(1.85), Inches(1.5))
        elif has_bullets:
            # 2行以下 → アイコンバレット
            _set(ph_map, 32, '')
            add_bullet_icons(slide, body, Inches(0.8), Inches(2.2), Inches(8.0))
            add_accent_line(slide, Inches(0.8), Inches(1.85), Inches(1.5))
        else:
            if 32 in ph_map:
                set_rich_body(ph_map[32], body)
            add_accent_line(slide, Inches(0.8), Inches(1.85), Inches(1.5))

    elif layout == 'two-column':
        _set(ph_map, 29, data.get('title', ''))
        if 30 in ph_map:
            set_rich_body(ph_map[30], data.get('leftBody', ''))
        if 32 in ph_map:
            set_rich_body(ph_map[32], data.get('rightBody', ''))
        _set(ph_map, 33, '')
        # 中央の区切り線
        add_accent_line(slide, Inches(4.85), Inches(2.0), Inches(0.04))

    elif layout == 'centered-text':
        _set(ph_map, 29, data.get('title', ''))
        _set(ph_map, 30, '')
        if 32 in ph_map:
            set_rich_body(ph_map[32], data.get('body', ''))
        elif 31 in ph_map:
            set_rich_body(ph_map[31], data.get('body', ''))

    elif layout in IMAGE_LAYOUTS:
        # 画像系レイアウト共通処理
        _set(ph_map, 29, data.get('title', ''))
        if 30 in ph_map:
            _set(ph_map, 30, '')
        body = data.get('body', '')
        if layout != 'picture-fullscreen':
            # テキスト+画像レイアウト
            if 32 in ph_map:
                set_rich_body(ph_map[32], body)
            elif 31 in ph_map:
                set_rich_body(ph_map[31], body)
        # 画像挿入は handler の do_POST 内で実行（image_map経由）
        add_accent_line(slide, Inches(0.8), Inches(1.85), Inches(1.5))

    elif layout == 'content-with-chart':
        _set(ph_map, 29, data.get('title', ''), font_size_pt=28)
        _set(ph_map, 30, '')
        body = data.get('body', '')
        if 32 in ph_map:
            set_rich_body(ph_map[32], body, base_size=14)
        # ネイティブチャート追加
        chart = data.get('chart')
        if chart and chart.get('labels') and chart.get('data'):
            add_chart(slide, chart)

    elif layout == 'content-with-flow':
        _set(ph_map, 29, data.get('title', ''), font_size_pt=28)
        _set(ph_map, 30, '')
        if 32 in ph_map:
            set_rich_body(ph_map[32], data.get('body', ''), base_size=14)
        # フロー図形追加
        flow = data.get('flow')
        if flow and flow.get('steps'):
            add_flow(slide, flow['steps'])

    elif layout == 'sixbox':
        _set(ph_map, 30, data.get('title', ''), font_size_pt=24)
        _set(ph_map, 31, '')
        boxes = data.get('boxes', [])
        for b in range(6):
            box = boxes[b] if b < len(boxes) else {}
            head_idx = (10 + b) if b < 3 else (11 + b)
            _set(ph_map, head_idx, box.get('heading', ''), font_size_pt=14)
            _set(ph_map, 17 + b, box.get('body', ''), font_size_pt=11)
            _set(ph_map, 23 + b, str(b + 1) if box.get('heading') else '')

    elif layout == 'comparison':
        _set(ph_map, 29, data.get('title', ''), font_size_pt=28)
        _set(ph_map, 30, '')
        _set(ph_map, 32, '')
        # ネイティブテーブル追加
        table = data.get('table')
        if table and table.get('headers'):
            add_table(slide, table)
        # タイトル下アクセントライン
        add_accent_line(slide, Inches(0.8), Inches(1.9), Inches(1.5))

    elif layout == 'quote':
        speaker = data.get('speaker', '')
        _set(ph_map, 30, f'- {speaker}' if speaker else '')
        _set(ph_map, 31, data.get('body', ''))
        # 大きな引用符マーク
        q_shape = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(1.0), Inches(0.8))
        qp = q_shape.text_frame.paragraphs[0]
        qp.text = '\u201C'  # "
        qp.font.size = Pt(72)
        qp.font.color.rgb = GOLD
        qp.font.bold = True

    elif layout == 'closing':
        msg = data.get('message', '')
        # Thank you テンプレートの利用可能なplaceholderに書き込む
        for idx in (29, 30, 31, 32, 0, 10):
            if idx in ph_map:
                try:
                    ph_map[idx].text = str(msg).strip() if msg else ''
                    break
                except Exception:
                    continue

    # スピーカーノートの追加（全レイアウト共通）
    notes_text = data.get('notes', '')
    if notes_text:
        try:
            notes_slide = slide.notes_slide
            tf = notes_slide.notes_text_frame
            tf.text = str(notes_text).strip()
        except Exception:
            pass


def build_layout_map(prs):
    layout_map = {}
    for master in prs.slide_masters:
        for sl in master.slide_layouts:
            layout_map[sl.name] = sl
    return layout_map


class handler(BaseHTTPRequestHandler):
    def do_POST(self):
        try:
            content_length = int(self.headers.get('Content-Length', 0))
            body = json.loads(self.rfile.read(content_length))

            slides_data = body.get('slides', [])
            template_id = body.get('template', 'external-white')

            if not slides_data:
                self._error(400, 'slides is required')
                return

            tpl_path = os.path.join(TEMPLATE_DIR, f'{template_id}.pptx')
            if not os.path.exists(tpl_path):
                self._error(400, f'Template not found: {template_id}')
                return

            prs = Presentation(tpl_path)
            layout_map = build_layout_map(prs)

            # 既存スライドを全削除
            sldIdLst = prs.slides._sldIdLst
            for sldId in list(sldIdLst):
                rId = sldId.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
                if rId:
                    prs.part.drop_rel(rId)
                sldIdLst.remove(sldId)

            # Unsplash未設定時: 画像系レイアウトをcontentにフォールバック
            if not UNSPLASH_ACCESS_KEY:
                for sd in slides_data:
                    if sd.get('layout') in IMAGE_LAYOUTS:
                        sd['_original_layout'] = sd['layout']
                        sd['layout'] = 'content'
                        # bodyが空の場合、imageQueryからヒントテキストを補完
                        if not sd.get('body') and sd.get('imageQuery'):
                            sd['body'] = sd.get('title', '')

            # 画像を並列取得
            image_map = fetch_images_batch(slides_data)

            for i, sd in enumerate(slides_data):
                ly = sd.get('layout', 'content')
                layout_name = LAYOUT_NAMES.get(ly, 'Title, Subtitle and one Paragrah')
                layout = layout_map.get(layout_name)
                if not layout:
                    layout = layout_map.get('Title, Subtitle and one Paragrah')
                if not layout:
                    continue

                slide = prs.slides.add_slide(layout)
                apply_data(slide, ly, sd)

                # 画像系レイアウト: PICTURE placeholderに画像挿入
                if ly in IMAGE_LAYOUTS and i in image_map:
                    # PICTURE placeholderのidxはレイアウトによって異なる
                    pic_idx = 12 if ly == 'content-left-full-picture-right' else (32 if ly == 'picture-fullscreen' else 21)
                    insert_picture_to_placeholder(slide, pic_idx, image_map[i])

            buffer = io.BytesIO()
            prs.save(buffer)
            pptx_b64 = base64.b64encode(buffer.getvalue()).decode('utf-8')

            self.send_response(200)
            self.send_header('Content-Type', 'application/json')
            self.send_header('Access-Control-Allow-Origin', '*')
            self.end_headers()
            self.wfile.write(json.dumps({
                'pptx': pptx_b64,
                'slides': len(slides_data),
            }).encode())

        except Exception as e:
            self._error(500, str(e))

    def do_OPTIONS(self):
        self.send_response(200)
        self.send_header('Access-Control-Allow-Origin', '*')
        self.send_header('Access-Control-Allow-Methods', 'POST, OPTIONS')
        self.send_header('Access-Control-Allow-Headers', 'Content-Type')
        self.end_headers()

    def _error(self, code, msg):
        self.send_response(code)
        self.send_header('Content-Type', 'application/json')
        self.send_header('Access-Control-Allow-Origin', '*')
        self.end_headers()
        self.wfile.write(json.dumps({'error': msg}).encode())
